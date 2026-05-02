from pptx import Presentation


def convert_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {slide_num}\n")

        title = slide.shapes.title
        if title and title.text.strip():
            parts.append(f"### {title.text.strip()}\n")

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == title:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    level = para.level
                    indent = "  " * level
                    parts.append(f"{indent}- {text}")
        parts.append("")

    return "\n".join(parts).strip()
